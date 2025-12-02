from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- Contenido de la presentación ---
slides_content = [
    {
        "title": "Python con Superpoderes: Una Introducción a los Type Hints",
        "subtitle": "Escribiendo código más claro, robusto y fácil de mantener.",
        "layout": "title_slide"
    },
    {
        "title": "¿Qué es Python? El Pato Escribe Como Puede",
        "points": [
            "Python es un lenguaje de tipado dinámico.",
            "No necesitas declarar el tipo de una variable; se infiere en tiempo de ejecución.",
            "Ventaja: Flexibilidad y rapidez al escribir.",
            "Desventaja: Puede llevar a errores inesperados."
        ],
        "code": '''variable = 5
print(type(variable))  # <class 'int'>

variable = "Hola Mundo"
print(type(variable))  # <class 'str'>'''
    },
    {
        "title": "¿Qué son los Type Hints?",
        "points": [
            "Son anotaciones (pistas) para indicar el tipo de dato esperado.",
            "Se aplican a variables, parámetros y valores de retorno de funciones.",
            "Introducidos oficialmente en PEP 484.",
            "¡OJO! Python ignora los hints en tiempo de ejecución. Son para desarrolladores y herramientas externas."
        ]
    },
    {
        "title": "¿Por qué usarlos? ¡Las Ventajas!",
        "points": [
            "Claridad y Autodocumentación: El código es más fácil de leer.",
            "Detección Temprana de Errores: Herramientas como `mypy` encuentran fallos antes de ejecutar.",
            "Mejor Soporte del Editor (IDE): Autocompletado más inteligente y refactorización segura.",
            "Mayor Robustez: Reduce errores en producción por tipos incorrectos."
        ]
    },
    {
        "title": "Sintaxis Básica: Variables y Funciones",
        "points": [
            ("Variables:", [": después del nombre de la variable."]),
            ("Funciones:", ["Parámetros: `:` después del nombre.", "Retorno: `->` antes de los dos puntos del `def`.", "Si no retorna nada, se usa `-> None`."])
        ],
        "code": '''# Variables
nombre: str = "Ana"
edad: int = 25

# Función que saluda
def saludar(nombre: str) -> str:
    return f"Hola, {nombre}"

# Función que no devuelve nada
def imprimir_mensaje(mensaje: str) -> None:
    print(mensaje)'''
    },
    {
        "title": "Tipos Complejos (Módulo typing)",
        "points": [
            "Para listas, diccionarios, etc., usamos el módulo `typing`.",
            "Listas: `list[int]` o `List[int]` (versiones antiguas).",
            "Diccionarios: `dict[str, int]` o `Dict[str, int]`."
        ],
        "code": '''from typing import List, Dict

# Sintaxis moderna (Python 3.9+)
def promediar(edades: list[int]) -> float:
    return sum(edades) / len(edades)

# Diccionario con tipos
puntuaciones: dict[str, int] = {"Ana": 90}'''
    },
    {
        "title": "Tipos Especiales del Módulo typing",
        "points": [
            "`Optional[str]` o `str | None`: El valor puede ser `str` o `None`.",
            "`Union[int, str]` o `int | str`: El valor puede ser de varios tipos.",
            "`Any`: El comodín. Para cualquier tipo. ¡Úsalo con moderación!"
        ],
        "code": '''from typing import Optional, Union, Any

def buscar_usuario(id: int) -> Optional[str]:
    if id == 1:
        return "Sergio"
    return None

def formatear_id(item_id: Union[int, str]) -> str:
    return f"ID-{item_id}"'''
    },
    {
        "title": "Verificación Estática con mypy",
        "points": [
            "Python no verifica los tipos; `mypy` sí.",
            "Instalación: `pip install mypy`.",
            "Uso: `mypy mi_codigo.py`."
        ],
        "code": '''# mi_codigo.py
def sumar(a: int, b: int) -> int:
    return a + b

sumar(5, "texto") # mypy mostrará un error aquí

# Salida de mypy:
# error: Argument 2 to "sumar" has incompatible 
# type "str"; expected "int"'''
    },
    {
        "title": "Conclusiones",
        "points": [
            "Los Type Hints mejoran la calidad de tu código.",
            "No afectan el rendimiento en ejecución.",
            "Facilitan el trabajo en equipo y el mantenimiento.",
            "Son una herramienta clave en el desarrollo Python moderno."
        ]
    },
    {
        "title": "¿Preguntas?",
        "subtitle": "",
        "layout": "title_slide"
    }
]

# --- Creación de la presentación ---
prs = Presentation()
# Usar un fondo oscuro y texto claro para un look de "código"
# Puedes comentar estas líneas si prefieres el fondo blanco por defecto
# prs.slide_width = Inches(16)
# prs.slide_height = Inches(9)

def create_slide(prs, content):
    layout_name = content.get("layout", "title_and_content")
    
    if layout_name == "title_slide":
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = content["title"]
        subtitle.text = content.get("subtitle", "")
    else:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = content["title"]
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()  # Limpiar contenido por defecto

        # Añadir puntos
        if "points" in content:
            for item in content["points"]:
                if isinstance(item, tuple):
                    p = tf.add_paragraph()
                    p.text = item[0]
                    p.font.bold = True
                    p.level = 0
                    for sub_item in item[1]:
                        p_sub = tf.add_paragraph()
                        p_sub.text = sub_item
                        p_sub.level = 1
                else:
                    p = tf.add_paragraph()
                    p.text = item
                    p.level = 0
        
        # Añadir bloque de código
        if "code" in content:
            left = Inches(5.5)
            top = Inches(1.8)
            width = Inches(8)
            height = Inches(5)
            
            if "points" not in content: # Si no hay puntos, el código ocupa más espacio
                left, top, width = Inches(1.5), Inches(1.8), Inches(11)

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf_code = txBox.text_frame
            tf_code.word_wrap = True
            
            p_code = tf_code.add_paragraph()
            p_code.text = content["code"]
            p_code.font.name = 'Courier New'
            p_code.font.size = Pt(14)
            p_code.font.color.rgb = RGBColor(0,0,0) # Texto negro

# Generar todas las diapositivas
for slide_content in slides_content:
    create_slide(prs, slide_content)

# Guardar la presentación
file_name = "presentacion_type_hints.pptx"
prs.save(file_name)

print(f"¡Presentación '{file_name}' creada con éxito!")
